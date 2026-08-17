from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

slides_data = [
    {
        "title": "Welcome to Chop with Rosty 🍽️",
        "image": "/home/ishola/.gemini/antigravity/brain/5e47defe-d30c-4470-8ae1-bc21d17e7d8c/chop_rosty_title_1785509516631.png",
        "text": "We've successfully transformed the underlying architecture and beautifully rebranded the entire application, complete with fresh UI elements and customized alerts."
    },
    {
        "title": "Phase 1: Robust Local Environment",
        "image": "/home/ishola/.gemini/antigravity/brain/5e47defe-d30c-4470-8ae1-bc21d17e7d8c/local_dev_env_1785509538917.png",
        "text": "- Setup local Supabase Docker containers for 0-latency testing\n- Migrated all authentication away from cloud dependency\n- Intercepted magic links locally via Inbucket to avoid API limits\n- Synced Prisma schema safely with the new local environment"
    },
    {
        "title": "Phase 2: Admin Dashboard & Customers",
        "image": "/home/ishola/.gemini/antigravity/brain/5e47defe-d30c-4470-8ae1-bc21d17e7d8c/admin_dashboard_ui_1785509554848.png",
        "text": "- Dynamic ADMIN role assignment synced between Supabase and Prisma\n- Revamped Customer List showing intuitive, auto-incrementing shortId\n- Polished Inventory stock panels with working Add/Edit Modal forms"
    },
    {
        "title": "Phase 3: Kitchen & Order Management",
        "image": "/home/ishola/.gemini/antigravity/brain/5e47defe-d30c-4470-8ae1-bc21d17e7d8c/kitchen_orders_1785509573346.png",
        "text": "- Dedicated full-page routes for each Order (/admin/orders/[id])\n- Powerful 'Edit Ingredients' functionality utilizing Prisma Transactions\n- Safely deduct or refund inventory stock in real-time when orders change\n- Track exact changes with a detailed OrderIngredientLog"
    },
    {
        "title": "What's Next: WhatsApp Alerts!",
        "image": "/home/ishola/.gemini/antigravity/brain/5e47defe-d30c-4470-8ae1-bc21d17e7d8c/whatsapp_integration_1785509666639.png",
        "text": "We are now ready to integrate Meta's WhatsApp Business API!\n\nThis will allow us to instantly send order confirmations, delivery updates, and low-stock alerts straight to phones."
    }
]

for data in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide layout with title
    
    title = slide.shapes.title
    title.text = data["title"]
    
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(5)
    if os.path.exists(data["image"]):
        slide.shapes.add_picture(data["image"], left, top, width=width)
        
    txBox = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(3), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = data["text"]
    p.font.size = Pt(18)

prs.save('/home/ishola/jar/compENG/sem-8/swe-project/Chop_with_Rosty_Progress.pptx')
print("PPTX generated successfully!")
